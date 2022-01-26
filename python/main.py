import self as self
import db
import os
from comtypes.client import CreateObject, Constants
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
import configparser
import yaml


def __init__(self):
    self.env = None
    self.config = None
    self.template = None
    self.slides = None


def read_env(self):
    env_file = os.path.normpath(os.path.join(os.path.dirname(__file__), 'files', self.template, 'env.ini'))
    self.env = configparser.ConfigParser()
    self.env.read(env_file)


def read_config(self):
    conf_file = os.path.normpath(os.path.join(os.path.dirname(__file__), 'files', self.template, 'config.yaml'))
    with open(conf_file) as f:
        self.config = yaml.safe_load(f)


def execute_query(self, query):
    if self.config and self.env.get('DEFAULT', 'db') == 'mysql':
        db_client = db.MySQLClient(self.env.get('DEFAULT', 'host'), self.env.get('DEFAULT', 'user'),
                                   self.env.get('DEFAULT', 'password'), self.env.get('DEFAULT', 'database'))
        return db_client.execute_query(query)


def create_table(table_slide, component):
    data = execute_query(self, component['query'])
    x, y, cx, cy = Inches(component['cords']['x']), Inches(component['cords']['y']), Inches(
        component['cords']['width']), Inches(component['cords']['height'])
    print(len(data))
    rows = len(data)
    columns = len(component['headers']) - 1
    table = table_slide.shapes.add_table(rows + 1, len(component['headers']), x,
                                         y, cx,
                                         cy).table
    cell = table.cell(0, 0)
    other_cell = table.cell(rows, 0)
    cell.merge(other_cell)
    table.cell(0, 0).vertical_anchor = MSO_ANCHOR.MIDDLE

    for i in range(0, len(component['headers'])):
        table.cell(0, i).text = component['headers'][i]
        table.cell(0, i).width = Inches(component['column_width'][i])

    for index, tuple_data in enumerate(data):
        for col in range(0, columns):
            if tuple_data[col] is not None:
                table.cell(index + 1, col + 1).text = str(tuple_data[col])
            else:
                table.cell(index + 1, col).text = ''


def create_chart(slide, component):
    if component["chartType"] == "COLUMN":
        create_column_chart(slide, component)
    if component["chartType"] == "BAR":
        create_bar_chart(slide, component)
    else:
        print("Unsupported Chart type!")


def create_column_chart(chart_slide, component):
    # define chart data ---------------------
    data = execute_query(self, component['query'])
    chart_data = CategoryChartData()
    categories = []
    series = []
    for index, tuple_data in enumerate(data):
        categories.append(tuple_data[0])
        series.append(tuple_data[1])

    chart_data.categories = categories
    chart_data.add_series('Series 1', tuple(series))
    x, y, cx, cy = Inches(component['cords']['x']), Inches(component['cords']['y']), Inches(
        component['cords']['width']), Inches(component['cords']['height'])

    chart = chart_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart
    set_chart_properties(chart, component)


def create_bar_chart(chart_slide, component):
    # define chart data ---------------------
    data = execute_query(self, component['query'])
    print(data)
    categories = []
    series = []
    for index, tuple_data in enumerate(data):
        categories.append(tuple_data[0])
        series.append(tuple_data[1])

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('Series 1', series)
    # add chart to slide --------------------
    x, y, cx, cy = Inches(component['cords']['x']), Inches(component['cords']['y']), Inches(
        component['cords']['width']), Inches(component['cords']['height'])
    chart = chart_slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED, x, y, cx, cy, chart_data
    ).chart
    set_chart_properties(chart, component)


def set_chart_properties(chart, component):
    plot = chart.plots[0]
    chart_series = plot.series[0]
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = False
    chart.chart_title.text_frame.text = component['title']
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(10)
    fill = chart_series.format.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(79, 129, 189)
    value_axis = chart.value_axis
    value_axis.minimum_scale = 0
    value_axis.maximum_scale = None
    value_axis.has_major_gridlines = False
    value_axis.has_minor_gridlines = False


def replace_text(replacements, slide_shapes):
    for shape in slide_shapes:
        if shape.has_text_frame:
            if (shape.text.find(replacements['key'])) != -1:
                text_frame = shape.text_frame
                match_and_replace(replacements, text_frame)
        elif shape.has_table:
            for cell in shape.table.iter_cells():
                if (cell.text.find(replacements['key'])) != -1:
                    text_frame = cell.text_frame
                    match_and_replace(replacements, text_frame)


def match_and_replace(replacements, text_frame):
    data = execute_query(self, replacements['query'])
    for paragraph in text_frame.paragraphs:
        whole_text = "".join(run.text for run in paragraph.runs)
        whole_text = whole_text.replace(str(replacements['key']), str(data[0][0]))
        for idx, run in enumerate(paragraph.runs):
            if idx != 0:
                p = paragraph._p
                p.remove(run._r)
        if not (not paragraph.runs):
            paragraph.runs[0].text = whole_text


def ppt_to_pdf(input_file_name, output_file_name, formatType=32):
    powerpoint = CreateObject('Powerpoint.Application')
    constants = Constants(powerpoint)
    powerpoint.Visible = 1
    if output_file_name[-3:] != 'pdf':
        output_file_name = output_file_name + ".pdf"
    deck = powerpoint.Presentations.Open(input_file_name)
    deck.SaveAs(output_file_name, constants.PpSaveAsPDF)
    deck.Close()
    powerpoint.Quit()


def export_as_pdf():
    mypath = os.path.abspath(__file__)
    mydir = os.path.dirname(mypath)
    file_input = os.path.join(mydir, self.config['fileName'] + ".pptx")
    # create the pdf output file path and call your function
    file_output = os.path.join(mydir, self.config['fileName'])
    ppt_to_pdf(file_input, file_output)


def validate_config_yaml():
    pass


if __name__ == '__main__':
    self.template = "templatename"
    read_config(self)
    validate_config_yaml()
    read_env(self)
    template_file_name = self.config['templateFileName']
    print(self.config['templateFileLocation'] + '\\' + template_file_name + '.pptx')

    prs = Presentation(self.config['templateFileLocation'] + '\\' + template_file_name + '.pptx')
    # To get shapes in your slides
    slides = [slide for slide in prs.slides]
    for slide in slides:
        if slide.shapes.title:
            print(slide.shapes.title)
            """ Filter  Slide object from Config yaml based on title name.
            Based on config slide component type, system 
            will either replace text or create table/chart. """
            result = list(filter(lambda x: (x['name'] == slide.shapes.title.text), self.config['slides']))
            if result:
                for item in result[0]['components']:
                    if item["type"] == "TEXT":
                        replace_text(item, slide.shapes)
                    elif item["type"] == "CHART":
                        create_chart(slide, item)
                    elif item["type"] == "TABLE":
                        create_table(slide, item)
            else:
                print("No Result found for given slide in config yaml!")

    prs.save(self.config['fileName'] + '.pptx')
    export_as_pdf()
